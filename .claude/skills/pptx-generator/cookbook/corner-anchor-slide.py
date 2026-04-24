#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "corner-anchor-slide"
# purpose = "Bold statement with strong visual anchor, asymmetric modern design"
# best_for = [
#     "Bold statements with supporting text",
#     "Modern/asymmetric visual style",
#     "Key points with emphasis",
#     "Visual variety in presentation",
# ]
# avoid_when = [
#     "Need lots of text or bullets",
#     "Already used this layout recently",
#     "Conservative/traditional audience",
# ]
# anchor_corners = ["top-left", "top-right", "bottom-left", "bottom-right"]
# instructions = [
#     "ANCHOR_CORNER determines where the large circle appears",
#     "Text is positioned opposite to the anchor for balance",
#     "Headline can be longer than title-slide (full statement)",
#     "Body text provides supporting context or details",
#     "Use sparingly for visual variety",
# ]
# ///
"""
LAYOUT: Corner Anchor
PURPOSE: Bold statement with strong visual anchor, asymmetric modern design

CUSTOMIZE:
- BRAND_* colors and fonts
- HEADLINE: Main statement
- BODY_TEXT: Supporting paragraph or points
- ANCHOR_CORNER: "top-left", "top-right", "bottom-left", "bottom-right"
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color (for main anchor shape)
    BRAND_ACCENT_SECONDARY = "REPLACE"  # secondary accent (for layered shape)
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    HEADLINE = "REPLACE"              # main statement
    BODY_TEXT = "REPLACE"             # supporting paragraph or points
    ANCHOR_CORNER = "REPLACE"         # "top-left", "top-right", "bottom-left", "bottom-right"

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large anchor shape based on corner
    anchor_size = 5.5
    if ANCHOR_CORNER == "top-right":
        # Large quarter circle in top-right
        anchor = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(13.333 - anchor_size/2), Inches(-anchor_size/2),
            Inches(anchor_size), Inches(anchor_size)
        )
        text_x, text_y = 0.8, 3.0
        text_align = PP_ALIGN.LEFT
    elif ANCHOR_CORNER == "top-left":
        anchor = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-anchor_size/2), Inches(-anchor_size/2),
            Inches(anchor_size), Inches(anchor_size)
        )
        text_x, text_y = 5.5, 3.0
        text_align = PP_ALIGN.RIGHT
    elif ANCHOR_CORNER == "bottom-right":
        anchor = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(13.333 - anchor_size/2), Inches(7.5 - anchor_size/2),
            Inches(anchor_size), Inches(anchor_size)
        )
        text_x, text_y = 0.8, 1.5
        text_align = PP_ALIGN.LEFT
    else:  # bottom-left
        anchor = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-anchor_size/2), Inches(7.5 - anchor_size/2),
            Inches(anchor_size), Inches(anchor_size)
        )
        text_x, text_y = 5.5, 1.5
        text_align = PP_ALIGN.RIGHT

    anchor.fill.solid()
    anchor.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    anchor.line.fill.background()

    # Secondary smaller shape for layering
    small_anchor = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(13.333 - 2.5), Inches(-0.5),
        Inches(2.5), Inches(2.5)
    )
    small_anchor.fill.solid()
    small_anchor.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY)
    small_anchor.line.fill.background()

    # Headline
    headline_box = slide.shapes.add_textbox(
        Inches(text_x), Inches(text_y),
        Inches(7.0), Inches(1.5)
    )
    tf = headline_box.text_frame
    p = tf.paragraphs[0]
    p.text = HEADLINE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = text_align

    # Accent line
    if text_align == PP_ALIGN.LEFT:
        line_x = text_x
    else:
        line_x = text_x + 7.0 - 2.5

    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(line_x), Inches(text_y + 1.6),
        Inches(2.5), Inches(0.08)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_line.line.fill.background()

    # Body text
    body_box = slide.shapes.add_textbox(
        Inches(text_x), Inches(text_y + 2.0),
        Inches(7.0), Inches(2.5)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = BODY_TEXT
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.alignment = text_align

    # Save
    output = Path("corner-anchor-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
