#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "image-caption-slide"
# purpose = "Screenshots, diagrams, photos with explanatory text"
# best_for = [
#     "Product screenshots",
#     "Architecture diagrams",
#     "Photos with context",
#     "Visual demonstrations",
# ]
# avoid_when = [
#     "No image available - use another layout",
#     "Image needs full slide - create custom full-bleed layout",
#     "Multiple images needed - consider a collage or multiple slides",
# ]
# instructions = [
#     "Image displays on right, text on left (adjustable in code)",
#     "If no image provided, shows styled placeholder",
#     "Caption should explain what to notice in the image",
#     "Image path should be absolute or relative to script location",
# ]
# ///
"""
LAYOUT: Image with Caption Slide
PURPOSE: Screenshots, diagrams, photos with explanatory text

CUSTOMIZE:
- BRAND_* colors and fonts
- IMAGE_PATH: Path to the image file (PNG, JPG, etc.)
- CAPTION: Explanation of what the image shows
- TITLE: Optional title above image
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    TITLE = "REPLACE"                 # title above image
    CAPTION = "REPLACE"               # explanation of what the image shows
    IMAGE_PATH = None                 # path to image file (PNG, JPG, etc.)

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Left side: Text content
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(5.5), Inches(1.2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = TITLE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Accent underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(2.8),
        Inches(2.0), Inches(0.06)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    underline.line.fill.background()

    # Caption text
    caption_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.2),
        Inches(5.5), Inches(3.0)
    )
    tf = caption_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = CAPTION
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.line_spacing = 1.4

    # Right side: Image placeholder
    # If no image provided, show a styled placeholder
    image_x = Inches(6.5)
    image_y = Inches(0.8)
    image_width = Inches(6.33)
    image_height = Inches(5.9)

    if IMAGE_PATH and Path(IMAGE_PATH).exists():
        slide.shapes.add_picture(
            str(IMAGE_PATH),
            image_x, image_y,
            width=image_width
        )
    else:
        # Placeholder box (styled to look intentional)
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            image_x, image_y,
            image_width, image_height
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = hex_to_rgb("181825")
        placeholder.line.color.rgb = hex_to_rgb("313244")

        # Placeholder icon (camera/image symbol made of shapes)
        # Mountain shape (simplified image icon)
        icon_center_x = 6.5 + (6.33 / 2)
        icon_center_y = 0.8 + (5.9 / 2)

        # Outer frame
        frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(icon_center_x - 1.0), Inches(icon_center_y - 0.7),
            Inches(2.0), Inches(1.4)
        )
        frame.fill.background()
        frame.line.color.rgb = hex_to_rgb(BRAND_ACCENT)
        frame.line.width = Pt(2)

        # "Image" text
        img_text = slide.shapes.add_textbox(
            Inches(icon_center_x - 1.0), Inches(icon_center_y - 0.2),
            Inches(2.0), Inches(0.5)
        )
        p = img_text.text_frame.paragraphs[0]
        p.text = "IMAGE"
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
        p.alignment = PP_ALIGN.CENTER

    # Decorative corner accent (top right)
    corner_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(12.833), Inches(0),
        Inches(0.5), Inches(0.08)
    )
    corner_bar.fill.solid()
    corner_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    corner_bar.line.fill.background()

    corner_bar_v = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(13.253), Inches(0),
        Inches(0.08), Inches(0.5)
    )
    corner_bar_v.fill.solid()
    corner_bar_v.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    corner_bar_v.line.fill.background()

    # Save
    output = Path("image-caption-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
