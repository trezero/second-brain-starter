#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "code-slide"
# purpose = "Display code snippets, terminal output, configuration examples"
# best_for = [
#     "Showing code examples",
#     "Terminal/CLI output",
#     "Configuration snippets",
#     "API request/response examples",
# ]
# avoid_when = [
#     "Code is longer than 12 lines - split into multiple slides",
#     "Code needs syntax highlighting (consider screenshot instead)",
#     "Audience won't understand code - use pseudocode in content-slide",
# ]
# max_lines = 12
# instructions = [
#     "Keep code to 5-12 lines - less is better for readability",
#     "Include language label badge (Python, TypeScript, etc.)",
#     "Use caption to explain what the code does or key insight",
#     "Remove unnecessary comments/whitespace to save space",
#     "If code is complex, highlight the important part only",
# ]
# ///
"""
LAYOUT: Code Slide
PURPOSE: Display code snippets, terminal output, configuration examples

CUSTOMIZE:
- BRAND_* colors and fonts
- TITLE: What the code demonstrates
- CODE_CONTENT: The actual code (keep it short, 5-12 lines ideal)
- LANGUAGE_LABEL: Optional language badge (e.g., "Python", "TypeScript")
- CAPTION: Optional explanation below code
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color
    BRAND_ACCENT_SECONDARY = "REPLACE"  # secondary accent (for code block)
    BRAND_CODE_BG = "REPLACE"         # darker color for code block bg
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name
    BRAND_CODE_FONT = "REPLACE"       # monospace font for code

    # === CONTENT (from user request) ===
    TITLE = "REPLACE"                 # what the code demonstrates
    LANGUAGE_LABEL = "REPLACE"        # e.g., "Python", "TypeScript"
    CODE_CONTENT = '''REPLACE'''      # the actual code (5-12 lines)
    CAPTION = "REPLACE"               # explanation below code

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4),
        Inches(10.0), Inches(0.8)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = TITLE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Language label badge (top right of code block)
    if LANGUAGE_LABEL:
        badge_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(11.5), Inches(1.15),
            Inches(1.3), Inches(0.4)
        )
        badge_bg.fill.solid()
        badge_bg.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
        badge_bg.line.fill.background()

        badge_text = slide.shapes.add_textbox(
            Inches(11.5), Inches(1.18),
            Inches(1.3), Inches(0.4)
        )
        p = badge_text.text_frame.paragraphs[0]
        p.text = LANGUAGE_LABEL
        p.font.name = BRAND_CODE_FONT
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_CODE_BG)
        p.alignment = PP_ALIGN.CENTER

    # Code block background (terminal-like)
    code_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.4),
        Inches(12.33), Inches(4.8)
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = hex_to_rgb(BRAND_CODE_BG)
    code_bg.line.color.rgb = hex_to_rgb("313244")  # Subtle border

    # Terminal chrome dots (top left of code block)
    dot_colors = ["f38ba8", "f9e2af", "a6e3a1"]  # Red, yellow, green
    for i, color in enumerate(dot_colors):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8 + (i * 0.35)), Inches(1.65),
            Inches(0.18), Inches(0.18)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb(color)
        dot.line.fill.background()

    # Code content
    code_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.2),
        Inches(11.7), Inches(3.8)
    )
    tf = code_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = CODE_CONTENT
    p.font.name = BRAND_CODE_FONT
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.LEFT

    # Caption below code
    if CAPTION:
        caption_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.4),
            Inches(12.33), Inches(0.6)
        )
        p = caption_box.text_frame.paragraphs[0]
        p.text = CAPTION
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.alignment = PP_ALIGN.CENTER

    # Accent bar at bottom
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.35),
        Inches(13.333), Inches(0.15)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY)
    bottom_bar.line.fill.background()

    # Save
    output = Path("code-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
