#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "title-slide"
# purpose = "Opening slide, section opener, or standalone statement"
# best_for = [
#     "Presentation opening",
#     "Section dividers",
#     "Single powerful statement",
#     "Workshop/talk titles",
# ]
# avoid_when = [
#     "Need bullet points or details",
#     "Showing data or comparisons",
# ]
# instructions = [
#     "Keep headline to 3-8 words for impact",
#     "Subtitle is optional - skip if headline is self-explanatory",
#     "This slide sets the tone - make the headline memorable",
# ]
# ///
"""
LAYOUT: Title Slide
PURPOSE: Opening slide, section opener, or standalone statement

CUSTOMIZE:
- BRAND_* colors and fonts
- HEADLINE: Main message (3-8 words, punchy)
- SUBTITLE: Supporting context (optional)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    HEADLINE = "REPLACE"              # main message (3-8 words)
    SUBTITLE = "REPLACE"              # supporting context

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Accent bar at top (subtle brand touch)
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    top_bar.line.fill.background()

    # Headline - centered, impactful
    headline_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.8),
        Inches(11.83), Inches(1.8)
    )
    tf = headline_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = HEADLINE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle - below headline
    if SUBTITLE:
        subtitle_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(4.6),
            Inches(10.33), Inches(1.0)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = SUBTITLE
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(24)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.alignment = PP_ALIGN.CENTER

    # Save
    output = Path("title-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
