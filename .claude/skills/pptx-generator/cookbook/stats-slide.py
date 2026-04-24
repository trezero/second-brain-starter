#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "stats-slide"
# purpose = "Big numbers, KPIs, impressive metrics, before/after comparisons"
# best_for = [
#     "Showcasing 2-4 key metrics",
#     "Before/after comparisons",
#     "ROI or impact numbers",
#     "Quick wins or achievements",
# ]
# avoid_when = [
#     "More than 4 statistics - pick the most impactful",
#     "Numbers need detailed context - use content-slide instead",
#     "Comparing complex datasets - use chart-slide instead",
# ]
# max_stats = 4
# instructions = [
#     "Values should be short: '47%', '3x', '$2.1M', not long sentences",
#     "Labels explain what the number means in 2-5 words",
#     "Title should frame WHY these numbers matter",
#     "Most impactful stat should come first",
# ]
# ///
"""
LAYOUT: Stats/Metrics Slide
PURPOSE: Big numbers, KPIs, impressive metrics, before/after comparisons

CUSTOMIZE:
- BRAND_* colors and fonts
- TITLE: What these numbers prove
- STATS: List of (value, label) tuples - 2-4 stats work best
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color
    BRAND_ACCENT_SECONDARY = "REPLACE"  # secondary accent (for variety)
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    TITLE = "REPLACE"                 # what these numbers prove
    STATS = [
        ("REPLACE", "REPLACE"),       # (value, label) - 2-4 stats
        ("REPLACE", "REPLACE"),
        ("REPLACE", "REPLACE"),
    ]

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Top accent line
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.3),
        Inches(3.0), Inches(0.06)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    top_line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.6),
        Inches(12.0), Inches(1.0)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = TITLE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Calculate spacing for stats
    num_stats = len(STATS)
    total_width = 12.0
    stat_width = total_width / num_stats
    start_x = 0.66

    # Alternate accent colors for visual interest
    accent_colors = [BRAND_ACCENT, BRAND_ACCENT_SECONDARY, "a6e3a1"]

    for i, (value, label) in enumerate(STATS):
        x_pos = start_x + (i * stat_width)
        accent_color = accent_colors[i % len(accent_colors)]

        # Value (big number)
        value_box = slide.shapes.add_textbox(
            Inches(x_pos), Inches(2.5),
            Inches(stat_width - 0.3), Inches(1.8)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(accent_color)
        p.alignment = PP_ALIGN.CENTER

        # Underline accent
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_pos + 0.5), Inches(4.4),
            Inches(stat_width - 1.3), Inches(0.06)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = hex_to_rgb(accent_color)
        underline.line.fill.background()

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(x_pos), Inches(4.6),
            Inches(stat_width - 0.3), Inches(1.0)
        )
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.alignment = PP_ALIGN.CENTER

    # Subtle dividers between stats (optional visual)
    for i in range(1, num_stats):
        x_pos = start_x + (i * stat_width) - 0.15
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_pos), Inches(2.8),
            Inches(0.02), Inches(2.5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = hex_to_rgb("313244")  # Subtle
        divider.line.fill.background()

    # Save
    output = Path("stats-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
