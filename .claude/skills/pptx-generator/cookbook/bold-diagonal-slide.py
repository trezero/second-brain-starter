#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "bold-diagonal-slide"
# purpose = "High-energy statements, disruption themes, modern/dynamic feel"
# best_for = [
#     "Bold proclamations",
#     "Disruption or change themes",
#     "High-energy moments",
#     "Startup/innovation vibes",
# ]
# avoid_when = [
#     "Content needs structure or bullets",
#     "Conservative/formal audience",
#     "Already used multiple dynamic layouts",
# ]
# instructions = [
#     "Headline should be punchy: 2-5 words maximum",
#     "Subtext provides context but keep it brief",
#     "Diagonal shapes create energy - use sparingly in presentation",
#     "Text is right-aligned to balance the left diagonal shapes",
# ]
# ///
"""
LAYOUT: Bold Diagonal
PURPOSE: High-energy statements, disruption themes, modern/dynamic feel

CUSTOMIZE:
- BRAND_* colors and fonts
- HEADLINE: Bold statement (2-5 words)
- SUBTEXT: Supporting context
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color (for primary diagonal)
    BRAND_ACCENT_SECONDARY = "REPLACE"  # secondary accent (for layered diagonal)
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    HEADLINE = "REPLACE"              # bold statement (2-5 words)
    SUBTEXT = "REPLACE"               # supporting context

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large diagonal accent block (using parallelogram)
    # Create diagonal effect with rotated freeform or multiple shapes
    # Since python-pptx doesn't support rotation easily, we fake it with overlapping shapes

    # Main diagonal block (bottom-left to top-right feel)
    # We'll create a series of rectangles to simulate diagonal
    diagonal_colors = [BRAND_ACCENT, BRAND_ACCENT_SECONDARY]

    # Large accent rectangle anchored to bottom-left, extending up
    block1 = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(-1), Inches(3),
        Inches(8), Inches(5)
    )
    block1.fill.solid()
    block1.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    block1.line.fill.background()

    # Second parallelogram for layered effect
    block2 = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(-2), Inches(4.5),
        Inches(7), Inches(4)
    )
    block2.fill.solid()
    block2.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY)
    block2.line.fill.background()

    # Headline - positioned in the clear area (right side)
    headline_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(2.0),
        Inches(7.3), Inches(2.0)
    )
    tf = headline_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = HEADLINE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.RIGHT

    # Subtext
    subtext_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(4.2),
        Inches(7.3), Inches(1.0)
    )
    p = subtext_box.text_frame.paragraphs[0]
    p.text = SUBTEXT
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(24)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.alignment = PP_ALIGN.RIGHT

    # Small accent detail - top right corner
    corner = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(12.333), Inches(0),
        Inches(1), Inches(1)
    )
    corner.fill.solid()
    corner.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    corner.line.fill.background()

    # Save
    output = Path("bold-diagonal-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
