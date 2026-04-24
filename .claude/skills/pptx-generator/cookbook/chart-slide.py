#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "chart-slide"
# purpose = "Data visualization, metrics breakdown, distribution charts"
# best_for = [
#     "Showing data distribution (pie/doughnut)",
#     "Comparing categories (bar/column)",
#     "Visualizing survey results",
#     "Budget or resource allocation",
# ]
# avoid_when = [
#     "More than 6 data categories - simplify first",
#     "Trend over time - consider line chart (not in template)",
#     "Exact numbers matter more than visual - use stats-slide",
# ]
# chart_types = ["pie", "bar", "column", "doughnut"]
# max_categories = 6
# instructions = [
#     "Chart types: 'pie', 'bar', 'column', or 'doughnut'",
#     "Keep to 3-6 categories for readability",
#     "INSIGHT box on right explains the key takeaway",
#     "Title should explain what the data represents",
#     "Chart uses PowerPoint default theme colors",
# ]
# ///
"""
LAYOUT: Chart Slide
PURPOSE: Data visualization, metrics breakdown, distribution charts

CUSTOMIZE:
- BRAND_* colors and fonts
- TITLE: What the data shows
- CHART_TYPE: "pie", "bar", "column", "doughnut"
- CHART_DATA: List of (category, value) tuples
- INSIGHT: Key takeaway from the data
"""

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    TITLE = "REPLACE"                 # what the data shows
    CHART_TYPE = "REPLACE"            # "pie", "bar", "column", or "doughnut"
    CHART_DATA = [
        ("REPLACE", 0),               # (category, value) tuples
        ("REPLACE", 0),
        ("REPLACE", 0),
        ("REPLACE", 0),
    ]
    INSIGHT = "REPLACE"               # key takeaway from the data

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4),
        Inches(12.333), Inches(0.8)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = TITLE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Accent underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(2.0), Inches(0.06)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    underline.line.fill.background()

    # Prepare chart data
    chart_data = CategoryChartData()
    chart_data.categories = [item[0] for item in CHART_DATA]
    chart_data.add_series("Values", [item[1] for item in CHART_DATA])

    # Determine chart type
    chart_types = {
        "pie": XL_CHART_TYPE.PIE,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    xl_chart_type = chart_types.get(CHART_TYPE, XL_CHART_TYPE.PIE)

    # Add chart
    chart_left = Inches(0.5)
    chart_top = Inches(1.8)
    chart_width = Inches(7.5)
    chart_height = Inches(5.0)

    chart_shape = slide.shapes.add_chart(
        xl_chart_type, chart_left, chart_top, chart_width, chart_height, chart_data
    )
    chart = chart_shape.chart

    # Style the chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    # Note: Full chart styling (colors, fonts) requires more complex XML manipulation
    # The chart will use PowerPoint's default theme colors

    # Insight box (right side)
    insight_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(2.5),
        Inches(4.3), Inches(2.5)
    )
    insight_bg.fill.solid()
    insight_bg.fill.fore_color.rgb = hex_to_rgb("313244")
    insight_bg.line.fill.background()

    # Insight label
    insight_label = slide.shapes.add_textbox(
        Inches(8.8), Inches(2.7),
        Inches(3.7), Inches(0.5)
    )
    p = insight_label.text_frame.paragraphs[0]
    p.text = "KEY INSIGHT"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    # Insight text
    insight_box = slide.shapes.add_textbox(
        Inches(8.8), Inches(3.2),
        Inches(3.7), Inches(1.5)
    )
    tf = insight_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = INSIGHT
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Bottom accent bar
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.35),
        Inches(13.333), Inches(0.15)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    bottom_bar.line.fill.background()

    # Save
    output = Path("chart-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
