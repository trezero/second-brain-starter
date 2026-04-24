#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "floating-cards-slide"
# purpose = "Feature highlights, process steps, multiple equal items with depth"
# best_for = [
#     "Exactly 3 related features or concepts",
#     "Process with 3 steps",
#     "Three pillars or principles",
#     "Visual depth and modern aesthetic",
# ]
# avoid_when = [
#     "More than 3 items - use multi-card-slide instead",
#     "Less than 3 items - use content-slide instead",
#     "Long card titles (over 15 characters)",
#     "Long descriptions needed",
# ]
# max_cards = 3
# card_title_max_chars = 15
# instructions = [
#     "EXACTLY 3 cards required - no more, no less",
#     "Card titles must be SHORT: 1-2 words, max 15 characters",
#     "Descriptions should be 1-2 short sentences",
#     "Cards overlap intentionally for depth - do not use for 4+ items",
#     "If titles are too long, abbreviate or use different layout",
# ]
# ///
"""
LAYOUT: Floating Cards
PURPOSE: Feature highlights, process steps, multiple equal items with depth

CUSTOMIZE:
- BRAND_* colors and fonts
- TITLE: What these cards represent
- CARDS: List of (title, description) tuples - EXACTLY 3 cards required
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_CARD_BG = "REPLACE"         # card background color
    BRAND_CARD_BG_ALT = "REPLACE"     # alternate card background (for variety)
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color (card 1)
    BRAND_ACCENT_SECONDARY = "REPLACE"  # secondary accent (card 2)
    BRAND_ACCENT_TERTIARY = "REPLACE"   # tertiary accent (card 3)
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    TITLE = "REPLACE"                 # what these cards represent
    CARDS = [
        ("REPLACE", "REPLACE"),       # (title, description) - 3 cards works best
        ("REPLACE", "REPLACE"),
        ("REPLACE", "REPLACE"),
    ]

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(12.333), Inches(1.0)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = TITLE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Card configurations - staggered positions for depth
    card_configs = [
        {"x": 0.8, "y": 2.0, "accent": BRAND_ACCENT, "bg": BRAND_CARD_BG},
        {"x": 4.5, "y": 2.4, "accent": BRAND_ACCENT_SECONDARY, "bg": BRAND_CARD_BG_ALT},
        {"x": 8.2, "y": 1.8, "accent": BRAND_ACCENT_TERTIARY, "bg": BRAND_CARD_BG},
    ]

    card_width = 4.2
    card_height = 4.0

    for i, (card_title, card_desc) in enumerate(CARDS):
        config = card_configs[i % len(card_configs)]
        x = config["x"]
        y = config["y"]
        accent = config["accent"]
        bg = config["bg"]

        # Card shadow (offset rectangle)
        shadow = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.15), Inches(y + 0.15),
            Inches(card_width), Inches(card_height)
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = hex_to_rgb("11111b")
        shadow.line.fill.background()

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(card_width), Inches(card_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(bg)
        card.line.fill.background()

        # Accent bar at top of card
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.3), Inches(y + 0.3),
            Inches(0.8), Inches(0.12)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = hex_to_rgb(accent)
        accent_bar.line.fill.background()

        # Card number
        num_box = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + 0.6),
            Inches(1.0), Inches(0.8)
        )
        p = num_box.text_frame.paragraphs[0]
        p.text = f"0{i + 1}"
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(accent)

        # Card title (with word wrap for longer titles)
        card_title_box = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + 1.5),
            Inches(card_width - 0.6), Inches(1.0)
        )
        tf = card_title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card_title
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

        # Card description
        card_desc_box = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + 2.6),
            Inches(card_width - 0.6), Inches(1.1)
        )
        tf = card_desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card_desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # Save
    output = Path("floating-cards-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
