#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "multi-card-slide"
# purpose = "Multiple items displayed as cards in a row, flexible 3-5 cards"
# best_for = [
#     "3-5 equal items (pillars, features, steps)",
#     "Process steps displayed horizontally",
#     "Feature comparison cards",
#     "Team members or product tiers",
# ]
# avoid_when = [
#     "Only 2 items - use two-column-slide instead",
#     "More than 5 items - split across multiple slides",
#     "Items need detailed descriptions - use content-slide",
# ]
# min_cards = 3
# max_cards = 5
# card_title_max_chars = 12
# instructions = [
#     "Supports 3, 4, or 5 cards - layout adjusts automatically",
#     "Card titles must be SHORT: 1-2 words max (12 characters)",
#     "Descriptions should be 1 short sentence",
#     "Cards are evenly spaced with no overlap",
#     "For longer titles, abbreviate or use content-slide",
# ]
# ///
"""
LAYOUT: Multi-Card Slide
PURPOSE: Multiple items displayed as cards in a row, flexible 3-5 cards

CUSTOMIZE:
- BRAND_* colors and fonts
- TITLE: What these cards represent
- CARDS: List of (title, description) tuples - 3 to 5 cards supported
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_CARD_BG = "REPLACE"         # card background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color
    BRAND_ACCENT_SECONDARY = "REPLACE"  # secondary accent
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    TITLE = "REPLACE"                 # what these cards represent
    CARDS = [
        ("REPLACE", "REPLACE"),       # (title, description) - 3 to 5 cards
        ("REPLACE", "REPLACE"),
        ("REPLACE", "REPLACE"),
        ("REPLACE", "REPLACE"),
        ("REPLACE", "REPLACE"),
    ]

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4),
        Inches(12.333), Inches(0.9)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = TITLE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Accent underline below title
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.35),
        Inches(2.5), Inches(0.06)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    underline.line.fill.background()

    # Calculate card dimensions based on number of cards
    num_cards = len(CARDS)
    margin = 0.5  # Left/right margin
    gap = 0.3     # Gap between cards
    total_width = 13.333 - (2 * margin) - ((num_cards - 1) * gap)
    card_width = total_width / num_cards
    card_height = 4.8
    start_y = 1.8

    # Alternate accent colors for variety
    accent_colors = [BRAND_ACCENT, BRAND_ACCENT_SECONDARY, BRAND_ACCENT, BRAND_ACCENT_SECONDARY, BRAND_ACCENT]

    for i, (card_title, card_desc) in enumerate(CARDS):
        x = margin + (i * (card_width + gap))
        y = start_y
        accent = accent_colors[i % len(accent_colors)]

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(card_width), Inches(card_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.fill.background()

        # Accent bar at top of card
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(card_width), Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = hex_to_rgb(accent)
        accent_bar.line.fill.background()

        # Card number
        num_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.3),
            Inches(0.8), Inches(0.6)
        )
        p = num_box.text_frame.paragraphs[0]
        p.text = f"0{i + 1}"
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(accent)

        # Card title (with word wrap)
        card_title_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 1.0),
            Inches(card_width - 0.4), Inches(1.2)
        )
        tf = card_title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card_title
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

        # Card description (with word wrap)
        card_desc_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 2.3),
            Inches(card_width - 0.4), Inches(2.2)
        )
        tf = card_desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card_desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # Save
    output = Path("multi-card-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
