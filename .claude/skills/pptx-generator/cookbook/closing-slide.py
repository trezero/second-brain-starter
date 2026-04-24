#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "closing-slide"
# purpose = "Final slide, call to action, contact info, thank you, next steps"
# best_for = [
#     "Presentation ending",
#     "Call to action",
#     "Contact information",
#     "Next steps or resources",
#     "Q&A slide",
# ]
# avoid_when = [
#     "Not the final slide - use title-slide instead",
#     "Need detailed content - this is for impact, not information",
# ]
# max_cta_items = 4
# instructions = [
#     "Headline should be actionable or memorable",
#     "CTA items are optional - use for next steps or resources",
#     "Subtext is for URL, email, or social handle",
#     "Keep it simple - this is the lasting impression",
# ]
# ///
"""
LAYOUT: Closing/CTA Slide
PURPOSE: Final slide, call to action, contact info, thank you, next steps

CUSTOMIZE:
- BRAND_* colors and fonts
- HEADLINE: The takeaway or CTA (e.g., "Start Building Today", "Questions?")
- SUBTEXT: Supporting info (e.g., URL, email, social handle)
- CTA_ITEMS: Optional list of next steps or resources
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_BG_ALT = "REPLACE"          # slightly darker variant for contrast
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    HEADLINE = "REPLACE"              # takeaway or CTA (e.g., "Start Building Today")
    SUBTEXT = "REPLACE"               # URL, email, or social handle
    CTA_ITEMS = [
        "REPLACE",                    # next steps or resources (2-4 items)
        "REPLACE",
        "REPLACE",
    ]

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large accent gradient-like block (bottom portion)
    bottom_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.0),
        Inches(13.333), Inches(2.5)
    )
    bottom_block.fill.solid()
    bottom_block.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
    bottom_block.line.fill.background()

    # Top accent bar
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.12)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    top_bar.line.fill.background()

    # Large headline - centered
    headline_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0),
        Inches(12.33), Inches(1.5)
    )
    tf = headline_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = HEADLINE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # Accent underline below headline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(3.6),
        Inches(2.33), Inches(0.08)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    underline.line.fill.background()

    # CTA Items (horizontal layout in bottom section)
    if CTA_ITEMS:
        num_items = len(CTA_ITEMS)
        item_width = 12.0 / num_items
        start_x = 0.66

        for i, item in enumerate(CTA_ITEMS):
            x_pos = start_x + (i * item_width)

            # Number badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos + (item_width / 2) - 0.25), Inches(5.3),
                Inches(0.5), Inches(0.5)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
            badge.line.fill.background()

            badge_text = slide.shapes.add_textbox(
                Inches(x_pos + (item_width / 2) - 0.25), Inches(5.35),
                Inches(0.5), Inches(0.5)
            )
            p = badge_text.text_frame.paragraphs[0]
            p.text = str(i + 1)
            p.font.name = BRAND_HEADING_FONT
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(BRAND_BG)
            p.alignment = PP_ALIGN.CENTER

            # Item text
            item_box = slide.shapes.add_textbox(
                Inches(x_pos), Inches(5.95),
                Inches(item_width), Inches(0.8)
            )
            tf = item_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item
            p.font.name = BRAND_BODY_FONT
            p.font.size = Pt(18)
            p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
            p.alignment = PP_ALIGN.CENTER

    # Subtext (URL/contact) - bottom center
    if SUBTEXT:
        subtext_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.0),
            Inches(12.33), Inches(0.4)
        )
        p = subtext_box.text_frame.paragraphs[0]
        p.text = SUBTEXT
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
        p.alignment = PP_ALIGN.CENTER

    # Save
    output = Path("closing-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
