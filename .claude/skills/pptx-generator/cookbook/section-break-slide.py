#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "section-break-slide"
# purpose = "Transition between major sections, pause points, chapter markers"
# best_for = [
#     "Dividing presentation into major parts",
#     "Signaling topic change",
#     "Creating visual breathing room",
#     "Numbered sections or chapters",
# ]
# avoid_when = [
#     "Presentation is short (under 10 slides)",
#     "No clear section divisions",
#     "Already using title-slide for divisions",
# ]
# instructions = [
#     "Section number is optional - use '01', '02' or 'Part 1', 'Part 2'",
#     "Keep section title to 2-5 words",
#     "Subtitle briefly previews what's coming",
#     "Use consistently throughout the presentation",
# ]
# ///
"""
LAYOUT: Section Break Slide
PURPOSE: Transition between major sections, pause points, chapter markers

CUSTOMIZE:
- BRAND_* colors and fonts
- SECTION_NUMBER: Optional numbering (e.g., "01", "Part 2")
- SECTION_TITLE: The section name (bold, large)
- SECTION_SUBTITLE: Brief description of what's coming (optional)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_BG_ALT = "REPLACE"          # slightly darker variant for contrast
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    SECTION_NUMBER = "REPLACE"        # e.g., "01", "Part 2"
    SECTION_TITLE = "REPLACE"         # section name
    SECTION_SUBTITLE = "REPLACE"      # brief description

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large accent block (left third of slide)
    accent_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(4.5), Inches(7.5)
    )
    accent_block.fill.solid()
    accent_block.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
    accent_block.line.fill.background()

    # Section number (large, in accent block)
    if SECTION_NUMBER:
        num_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8),
            Inches(3.5), Inches(1.5)
        )
        p = num_box.text_frame.paragraphs[0]
        p.text = SECTION_NUMBER
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(120)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
        p.alignment = PP_ALIGN.CENTER

    # Vertical accent line (divider)
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.5), Inches(1.5),
        Inches(0.08), Inches(4.5)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    divider.line.fill.background()

    # Section title (right side)
    title_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(2.6),
        Inches(7.5), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = SECTION_TITLE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.LEFT

    # Section subtitle (below title)
    if SECTION_SUBTITLE:
        subtitle_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(4.2),
            Inches(7.5), Inches(1.0)
        )
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = SECTION_SUBTITLE
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.alignment = PP_ALIGN.LEFT

    # Save
    output = Path("section-break-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
