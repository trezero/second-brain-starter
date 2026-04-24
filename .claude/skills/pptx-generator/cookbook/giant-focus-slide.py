#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "giant-focus-slide"
# purpose = "Single powerful word/number, dramatic emphasis, memorable moments"
# best_for = [
#     "Single impactful word or number",
#     "Dramatic reveal moments",
#     "Key statistics (e.g., '47%', '$1B')",
#     "Memorable one-word themes",
# ]
# avoid_when = [
#     "Need to explain context in detail",
#     "More than 3 words needed",
#     "Already used this layout recently",
# ]
# max_words = 3
# instructions = [
#     "BIG_TEXT should be 1-3 words maximum - ideally just ONE word",
#     "Context above/below are optional but help frame the big text",
#     "Use for dramatic effect - sparingly in a presentation",
#     "Best for numbers, single concepts, or emotional words",
# ]
# ///
"""
LAYOUT: Giant Focus
PURPOSE: Single powerful word/number, dramatic emphasis, memorable moments

CUSTOMIZE:
- BRAND_* colors and fonts
- BIG_TEXT: The single word or short phrase (1-3 words max)
- CONTEXT_ABOVE: Small text above the big text
- CONTEXT_BELOW: Small text below the big text
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color (for the giant text)
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    BIG_TEXT = "REPLACE"              # single word or short phrase (1-3 words max)
    CONTEXT_ABOVE = "REPLACE"         # small text above the big text
    CONTEXT_BELOW = "REPLACE"         # small text below the big text

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large background accent shape (subtle)
    bg_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(3.5), Inches(0.5),
        Inches(6.333), Inches(6.333)
    )
    bg_circle.fill.solid()
    bg_circle.fill.fore_color.rgb = hex_to_rgb("232333")  # Very subtle
    bg_circle.line.fill.background()

    # Context above
    if CONTEXT_ABOVE:
        above_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8),
            Inches(12.333), Inches(0.6)
        )
        p = above_box.text_frame.paragraphs[0]
        p.text = CONTEXT_ABOVE.upper()
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.alignment = PP_ALIGN.CENTER
        # Letter spacing would be nice but not supported

    # GIANT TEXT
    big_box = slide.shapes.add_textbox(
        Inches(0), Inches(2.2),
        Inches(13.333), Inches(3.5)
    )
    tf = big_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = BIG_TEXT
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(200)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    p.alignment = PP_ALIGN.CENTER

    # Context below
    if CONTEXT_BELOW:
        below_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.6),
            Inches(12.333), Inches(0.6)
        )
        p = below_box.text_frame.paragraphs[0]
        p.text = CONTEXT_BELOW
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
        p.alignment = PP_ALIGN.CENTER

    # Accent line below context
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(6.4),
        Inches(2.333), Inches(0.04)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_line.line.fill.background()

    # Save
    output = Path("giant-focus-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
