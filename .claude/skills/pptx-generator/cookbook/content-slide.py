#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "content-slide"
# purpose = "Key points, takeaways, feature lists, explanations"
# best_for = [
#     "Listing 3-6 related points",
#     "Feature lists",
#     "Takeaways or summaries",
#     "Step-by-step instructions",
# ]
# avoid_when = [
#     "More than 6 bullet points - split into multiple slides",
#     "Need visual comparisons",
#     "Single big idea - use title-slide or giant-focus instead",
# ]
# max_bullets = 6
# instructions = [
#     "Keep bullets parallel in structure (all start with verbs, or all noun phrases)",
#     "Title should describe WHAT these points prove, not 'Key Points'",
#     "Each bullet should be 1-2 lines max",
#     "No trailing periods on bullets unless full sentences",
# ]
# ///
"""
LAYOUT: Content Slide
PURPOSE: Key points, takeaways, feature lists, explanations

CUSTOMIZE:
- BRAND_* colors and fonts
- TITLE: What this slide proves or explains (not "Key Points")
- BULLETS: 3-6 points, parallel structure, action-oriented
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color
    BRAND_ACCENT_SECONDARY = "REPLACE"  # secondary accent color
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    TITLE = "REPLACE"                 # slide title
    BULLETS = [
        "REPLACE",                    # 3-6 bullet points
        "REPLACE",
        "REPLACE",
    ]

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Left accent bar (vertical stripe)
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.08), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    left_bar.line.fill.background()

    # Title with accent underline
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.5),
        Inches(12.0), Inches(1.0)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = TITLE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Accent underline below title
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.5),
        Inches(2.5), Inches(0.06)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY)
    underline.line.fill.background()

    # Bullets with custom styling
    bullet_start_y = 2.0
    bullet_spacing = 1.1

    for i, bullet_text in enumerate(BULLETS):
        y_pos = bullet_start_y + (i * bullet_spacing)

        # Bullet indicator (small accent square)
        indicator = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(y_pos + 0.15),
            Inches(0.12), Inches(0.12)
        )
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
        indicator.line.fill.background()

        # Bullet text
        text_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(y_pos),
            Inches(11.5), Inches(0.9)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet_text
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(22)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Slide number indicator (bottom right)
    num_box = slide.shapes.add_textbox(
        Inches(12.5), Inches(7.0),
        Inches(0.5), Inches(0.3)
    )
    p = num_box.text_frame.paragraphs[0]
    p.text = "02"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.alignment = PP_ALIGN.RIGHT

    # Save
    output = Path("content-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
