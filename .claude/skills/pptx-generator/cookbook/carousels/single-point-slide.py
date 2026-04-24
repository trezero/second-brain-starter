#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "carousel-single-point-slide"
# format = "carousel"
# dimensions = "square"
# purpose = "One key point with supporting text - the workhorse slide"
# best_for = [
#     "Main content slides in carousel body",
#     "Single concept that needs explanation",
#     "Key takeaways or insights",
#     "Standalone valuable points",
# ]
# avoid_when = [
#     "First slide - use hook-slide instead",
#     "Last slide - use cta-slide instead",
#     "Numbered lists - use numbered-point-slide instead",
#     "Quotes - use quote-slide instead",
# ]
# headline_max_chars = 40
# body_max_chars = 150
# instructions = [
#     "Headline should be the KEY POINT in 40 chars or less",
#     "Body text expands on the headline - max 150 chars",
#     "One idea per slide - if you have multiple, split them",
#     "Make headline readable at thumbnail size",
# ]
# ///
"""
LAYOUT: Carousel Single Point Slide
PURPOSE: One key point with supporting explanation
FORMAT: Square (1:1) for LinkedIn carousels

CUSTOMIZE:
- BRAND_* colors and fonts
- HEADLINE: The main point (max 40 chars)
- BODY_TEXT: Supporting explanation (max 150 chars)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS ===
    BRAND_BG = "1e1e2e"
    BRAND_TEXT = "cdd6f4"
    BRAND_TEXT_SECONDARY = "bac2de"
    BRAND_ACCENT = "fab387"
    BRAND_HEADING_FONT = "JetBrains Mono"
    BRAND_BODY_FONT = "Inter"

    # === CONTENT ===
    HEADLINE = "REPLACE"  # Max 40 chars
    BODY_TEXT = "REPLACE"  # Max 150 chars

    # === SLIDE SETUP (Square 1:1) ===
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Left accent bar
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.1), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    left_bar.line.fill.background()

    # Headline
    headline_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0),
        Inches(6.5), Inches(2.0)
    )
    tf = headline_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = HEADLINE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Accent underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(4.1),
        Inches(2.0), Inches(0.06)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    underline.line.fill.background()

    # Body text
    body_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5),
        Inches(6.5), Inches(2.5)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = BODY_TEXT
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(24)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # Save
    output = Path("carousel-single-point-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
