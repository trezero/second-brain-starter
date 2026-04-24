#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "carousel-quote-slide"
# format = "carousel"
# dimensions = "square"
# purpose = "Quote with attribution - for social proof or memorable statements"
# best_for = [
#     "Social proof from customers or experts",
#     "Memorable statements worth highlighting",
#     "Your own quotable insights",
#     "Breaking up content with a different format",
# ]
# avoid_when = [
#     "Long quotes over 120 characters - paraphrase instead",
#     "Quotes that need context - add context in previous slide",
#     "First slide - use hook-slide instead",
# ]
# quote_max_chars = 120
# attribution_max_chars = 40
# instructions = [
#     "Quote should be IMPACTFUL and max 120 characters",
#     "Attribution is optional but recommended for credibility",
#     "Use quotation marks in the design, not the text",
#     "If it's your own quote, attribution can be your handle",
# ]
# ///
"""
LAYOUT: Carousel Quote Slide
PURPOSE: Featured quote with attribution
FORMAT: Square (1:1) for LinkedIn carousels

CUSTOMIZE:
- BRAND_* colors and fonts
- QUOTE_TEXT: The quote (max 120 chars, no quotation marks)
- ATTRIBUTION: Who said it (max 40 chars, optional)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS ===
    BRAND_BG = "1e1e2e"
    BRAND_TEXT = "cdd6f4"
    BRAND_TEXT_SECONDARY = "bac2de"
    BRAND_ACCENT = "fab387"
    BRAND_ACCENT_SECONDARY = "89b4fa"
    BRAND_HEADING_FONT = "JetBrains Mono"
    BRAND_BODY_FONT = "Inter"

    # === CONTENT ===
    QUOTE_TEXT = "REPLACE"  # Max 120 chars, no quotation marks
    ATTRIBUTION = ""  # Optional, max 40 chars

    # === SLIDE SETUP (Square 1:1) ===
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large opening quote mark
    open_quote = slide.shapes.add_textbox(
        Inches(0.3), Inches(1.0),
        Inches(2.0), Inches(2.0)
    )
    p = open_quote.text_frame.paragraphs[0]
    p.text = """
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(180)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    # Quote text
    quote_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(6.5), Inches(3.0)
    )
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = QUOTE_TEXT
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # Attribution
    if ATTRIBUTION:
        attr_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.8),
            Inches(6.5), Inches(0.8)
        )
        tf = attr_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"— {ATTRIBUTION}"
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(18)
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY)
        p.alignment = PP_ALIGN.CENTER

    # Bottom accent line
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.5), Inches(6.8),
        Inches(2.5), Inches(0.06)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    bottom_line.line.fill.background()

    # Save
    output = Path("carousel-quote-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
