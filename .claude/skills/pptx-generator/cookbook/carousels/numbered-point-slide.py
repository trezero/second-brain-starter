#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "carousel-numbered-point-slide"
# format = "carousel"
# dimensions = "square"
# purpose = "Numbered list item with big number - for listicle carousels"
# best_for = [
#     "Listicle content (5 tips, 7 mistakes, etc.)",
#     "Step-by-step processes",
#     "Ranked items",
#     "Sequential content where order matters",
# ]
# avoid_when = [
#     "Standalone points without sequence - use single-point-slide",
#     "First slide - use hook-slide with the list premise",
#     "Last slide - use cta-slide",
# ]
# point_text_max_chars = 60
# supporting_max_chars = 100
# instructions = [
#     "NUMBER should be the position (1, 2, 3, etc.)",
#     "POINT_TEXT is the main content - max 60 chars",
#     "SUPPORTING_TEXT is optional explanation - max 100 chars",
#     "Keep consistent numbering style across carousel",
# ]
# ///
"""
LAYOUT: Carousel Numbered Point Slide
PURPOSE: Single numbered item from a list - big number with point
FORMAT: Square (1:1) for LinkedIn carousels

CUSTOMIZE:
- BRAND_* colors and fonts
- NUMBER: The list position (1, 2, 3, etc.)
- POINT_TEXT: The main point (max 60 chars)
- SUPPORTING_TEXT: Optional explanation (max 100 chars)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS ===
    BRAND_BG = "1e1e2e"
    BRAND_BG_ALT = "181825"
    BRAND_TEXT = "cdd6f4"
    BRAND_TEXT_SECONDARY = "bac2de"
    BRAND_ACCENT = "fab387"
    BRAND_HEADING_FONT = "JetBrains Mono"
    BRAND_BODY_FONT = "Inter"

    # === CONTENT ===
    NUMBER = "1"  # Position in list
    POINT_TEXT = "REPLACE"  # Max 60 chars
    SUPPORTING_TEXT = ""  # Optional, max 100 chars

    # === SLIDE SETUP (Square 1:1) ===
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Number block (left side)
    num_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(2.5), Inches(7.5)
    )
    num_block.fill.solid()
    num_block.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
    num_block.line.fill.background()

    # Big number
    num_box = slide.shapes.add_textbox(
        Inches(0), Inches(2.5),
        Inches(2.5), Inches(2.5)
    )
    p = num_box.text_frame.paragraphs[0]
    p.text = NUMBER
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    p.alignment = PP_ALIGN.CENTER

    # Accent divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.5), Inches(1.5),
        Inches(0.08), Inches(4.5)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    divider.line.fill.background()

    # Point text
    point_box = slide.shapes.add_textbox(
        Inches(3.0), Inches(2.2),
        Inches(4.0), Inches(2.5)
    )
    tf = point_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = POINT_TEXT
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Supporting text
    if SUPPORTING_TEXT:
        support_box = slide.shapes.add_textbox(
            Inches(3.0), Inches(4.8),
            Inches(4.0), Inches(2.0)
        )
        tf = support_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = SUPPORTING_TEXT
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # Save
    output = Path("carousel-numbered-point-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
