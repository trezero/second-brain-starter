#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "carousel-hook-slide"
# format = "carousel"
# dimensions = "square"
# purpose = "Opening slide to grab attention - the scroll-stopper"
# best_for = [
#     "First slide of any carousel",
#     "Bold provocative statements",
#     "Questions that create curiosity",
#     "Controversial or surprising claims",
# ]
# avoid_when = [
#     "Middle of carousel - use single-point-slide instead",
#     "Long explanatory text - save that for body slides",
#     "Multiple points - this is ONE hook only",
# ]
# hook_max_chars = 50
# instructions = [
#     "Hook text should be SHORT and PUNCHY - max 50 characters",
#     "Use confrontational or curiosity-driven language",
#     "Subtext is optional - only if it adds value",
#     "This slide makes people stop scrolling",
# ]
# ///
"""
LAYOUT: Carousel Hook Slide
PURPOSE: Opening slide that stops the scroll - bold statement or question
FORMAT: Square (1:1) for LinkedIn carousels

CUSTOMIZE:
- BRAND_* colors and fonts
- HOOK_TEXT: The attention-grabbing statement (max 50 chars)
- SUBTEXT: Optional supporting line
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS ===
    BRAND_BG = "1e1e2e"
    BRAND_TEXT = "cdd6f4"
    BRAND_TEXT_SECONDARY = "bac2de"
    BRAND_ACCENT = "fab387"
    BRAND_HEADING_FONT = "JetBrains Mono"
    BRAND_BODY_FONT = "Inter"

    # === CONTENT ===
    HOOK_TEXT = "REPLACE"  # Max 50 chars, punchy
    SUBTEXT = ""  # Optional

    # === SLIDE SETUP (Square 1:1) ===
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Top accent bar
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(7.5), Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    top_bar.line.fill.background()

    # Large hook text - centered vertically
    hook_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(6.5), Inches(2.5)
    )
    tf = hook_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = HOOK_TEXT
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # Optional subtext
    if SUBTEXT:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.2),
            Inches(6.5), Inches(1.0)
        )
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = SUBTEXT
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(24)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.alignment = PP_ALIGN.CENTER

    # Bottom accent line
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.5), Inches(6.8),
        Inches(2.5), Inches(0.06)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    bottom_line.line.fill.background()

    # Save
    output = Path("carousel-hook-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
