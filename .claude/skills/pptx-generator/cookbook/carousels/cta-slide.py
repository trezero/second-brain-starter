#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "carousel-cta-slide"
# format = "carousel"
# dimensions = "square"
# purpose = "Closing slide with call to action - the conversion slide"
# best_for = [
#     "Last slide of any carousel",
#     "Follow/subscribe prompts",
#     "Link to resource or content",
#     "Summary with next step",
# ]
# avoid_when = [
#     "First or middle slides - save CTA for the end",
#     "Multiple CTAs - pick ONE action",
#     "No clear action - every carousel needs a CTA",
# ]
# headline_max_chars = 30
# cta_max_chars = 25
# handle_max_chars = 20
# instructions = [
#     "HEADLINE should summarize value or prompt action - max 30 chars",
#     "CTA_TEXT is the specific action - max 25 chars (e.g., 'Follow for more')",
#     "HANDLE is your social handle - max 20 chars",
#     "Keep it simple - one clear action beats multiple options",
# ]
# ///
"""
LAYOUT: Carousel CTA Slide
PURPOSE: Closing slide with call to action
FORMAT: Square (1:1) for LinkedIn carousels

CUSTOMIZE:
- BRAND_* colors and fonts
- HEADLINE: Summary or prompt (max 30 chars)
- CTA_TEXT: The action (max 25 chars)
- HANDLE: Your social handle (max 20 chars)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS ===
    BRAND_BG = "1e1e2e"
    BRAND_BG_ALT = "181825"
    BRAND_TEXT = "cdd6f4"
    BRAND_TEXT_SECONDARY = "bac2de"
    BRAND_ACCENT = "fab387"
    BRAND_HEADING_FONT = "JetBrains Mono"
    BRAND_BODY_FONT = "Inter"

    # === CONTENT ===
    HEADLINE = "REPLACE"  # Max 30 chars
    CTA_TEXT = "REPLACE"  # Max 25 chars
    HANDLE = "@YourHandle"  # Max 20 chars

    # === SLIDE SETUP (Square 1:1) ===
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Top accent bar
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(7.5), Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    top_bar.line.fill.background()

    # Bottom block
    bottom_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.0),
        Inches(7.5), Inches(2.5)
    )
    bottom_block.fill.solid()
    bottom_block.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
    bottom_block.line.fill.background()

    # Headline
    headline_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0),
        Inches(6.5), Inches(2.0)
    )
    tf = headline_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = HEADLINE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # CTA text
    cta_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.3),
        Inches(6.5), Inches(1.0)
    )
    tf = cta_box.text_frame
    p = tf.paragraphs[0]
    p.text = CTA_TEXT
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    p.alignment = PP_ALIGN.CENTER

    # Handle
    handle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.3),
        Inches(6.5), Inches(0.8)
    )
    tf = handle_box.text_frame
    p = tf.paragraphs[0]
    p.text = HANDLE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.alignment = PP_ALIGN.CENTER

    # Save
    output = Path("carousel-cta-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
