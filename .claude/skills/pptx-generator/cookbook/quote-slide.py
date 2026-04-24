#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "quote-slide"
# purpose = "Testimonials, pull quotes, key statements, memorable phrases"
# best_for = [
#     "Customer testimonials",
#     "Expert quotes",
#     "Memorable statements to emphasize",
#     "Social proof",
# ]
# avoid_when = [
#     "Quote is longer than 3 sentences",
#     "No clear attribution source",
#     "Need to show data alongside the quote",
# ]
# instructions = [
#     "Keep quotes to 1-3 sentences maximum",
#     "Attribution format: '— Name, Title at Company'",
#     "Choose quotes that support your narrative",
#     "If quote is too long, extract the most powerful part",
# ]
# ///
"""
LAYOUT: Quote Slide
PURPOSE: Testimonials, pull quotes, key statements, memorable phrases

CUSTOMIZE:
- BRAND_* colors and fonts
- QUOTE: The quote text (keep it punchy, 1-3 sentences max)
- ATTRIBUTION: Who said it (name, title, company)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    QUOTE = "REPLACE"                 # the quote text (1-3 sentences)
    ATTRIBUTION = "REPLACE"           # who said it (— Name, Title)

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large quote mark (decorative)
    quote_mark = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(2.0), Inches(2.0)
    )
    p = quote_mark.text_frame.paragraphs[0]
    p.text = "\u201C"  # Left double quote mark
    p.font.name = "Georgia"
    p.font.size = Pt(200)
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    # Make it semi-transparent by using a muted version
    # (python-pptx doesn't support opacity, so we use color)

    # Left accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0), Inches(2.8),
        Inches(0.12), Inches(2.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_bar.line.fill.background()

    # Quote text
    quote_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.8),
        Inches(10.5), Inches(2.5)
    )
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = QUOTE
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(36)
    p.font.italic = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.LEFT

    # Attribution
    attr_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(5.5),
        Inches(10.5), Inches(0.6)
    )
    p = attr_box.text_frame.paragraphs[0]
    p.text = ATTRIBUTION
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    p.alignment = PP_ALIGN.LEFT

    # Save
    output = Path("quote-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
