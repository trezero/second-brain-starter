#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "two-column-slide"
# purpose = "Before/after, pros/cons, old vs new, A vs B comparisons"
# best_for = [
#     "Before/after comparisons",
#     "Pros vs cons",
#     "Problem vs solution",
#     "Old way vs new way",
#     "Two contrasting approaches",
# ]
# avoid_when = [
#     "Comparing more than 2 things - consider a table or multiple slides",
#     "Items aren't truly comparable",
#     "One side has much more content than the other",
# ]
# max_bullets_per_column = 4
# instructions = [
#     "Left column is typically 'before/problem/old' (muted or red-ish color)",
#     "Right column is typically 'after/solution/new' (accent or green-ish color)",
#     "Keep bullet counts balanced between columns",
#     "Use color coding to reinforce the contrast",
# ]
# ///
"""
LAYOUT: Two-Column Comparison
PURPOSE: Before/after, pros/cons, old vs new, A vs B comparisons

CUSTOMIZE:
- BRAND_* colors and fonts
- TITLE: Frames the comparison
- LEFT_HEADER, LEFT_BULLETS: The "before" or "problem" side
- RIGHT_HEADER, RIGHT_BULLETS: The "after" or "solution" side
- LEFT_COLOR, RIGHT_COLOR: Color coding (e.g., red vs green, muted vs accent)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === COLUMN COLORS (choose contrasting colors for comparison) ===
    LEFT_COLOR = "REPLACE"            # problem/before side (e.g., red-ish)
    RIGHT_COLOR = "REPLACE"           # solution/after side (e.g., green-ish)

    # === CONTENT (from user request, written in brand voice) ===
    TITLE = "REPLACE"                 # frames the comparison
    LEFT_HEADER = "REPLACE"           # left column header
    LEFT_BULLETS = [
        "REPLACE",                    # left column points
        "REPLACE",
        "REPLACE",
    ]
    RIGHT_HEADER = "REPLACE"          # right column header
    RIGHT_BULLETS = [
        "REPLACE",                    # right column points
        "REPLACE",
        "REPLACE",
    ]

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Title bar background
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.3)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = hex_to_rgb("181825")  # Slightly darker
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4),
        Inches(12.33), Inches(0.8)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = TITLE
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Center divider (vertical line with glow effect)
    divider_glow = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.55), Inches(1.5),
        Inches(0.23), Inches(5.5)
    )
    divider_glow.fill.solid()
    divider_glow.fill.fore_color.rgb = hex_to_rgb("313244")
    divider_glow.line.fill.background()

    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.63), Inches(1.5),
        Inches(0.07), Inches(5.5)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    divider.line.fill.background()

    # LEFT SIDE
    # Header with color indicator
    left_indicator = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.7),
        Inches(0.15), Inches(0.5)
    )
    left_indicator.fill.solid()
    left_indicator.fill.fore_color.rgb = hex_to_rgb(LEFT_COLOR)
    left_indicator.line.fill.background()

    left_header_box = slide.shapes.add_textbox(
        Inches(0.85), Inches(1.65),
        Inches(5.5), Inches(0.6)
    )
    p = left_header_box.text_frame.paragraphs[0]
    p.text = LEFT_HEADER
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(LEFT_COLOR)

    # Left bullets
    for i, bullet in enumerate(LEFT_BULLETS):
        y_pos = 2.5 + (i * 1.2)

        # Bullet marker
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5), Inches(y_pos + 0.12),
            Inches(0.1), Inches(0.1)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = hex_to_rgb(LEFT_COLOR)
        marker.line.fill.background()

        # Bullet text
        text_box = slide.shapes.add_textbox(
            Inches(0.85), Inches(y_pos),
            Inches(5.5), Inches(1.0)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # RIGHT SIDE
    # Header with color indicator
    right_indicator = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.93), Inches(1.7),
        Inches(0.15), Inches(0.5)
    )
    right_indicator.fill.solid()
    right_indicator.fill.fore_color.rgb = hex_to_rgb(RIGHT_COLOR)
    right_indicator.line.fill.background()

    right_header_box = slide.shapes.add_textbox(
        Inches(7.28), Inches(1.65),
        Inches(5.5), Inches(0.6)
    )
    p = right_header_box.text_frame.paragraphs[0]
    p.text = RIGHT_HEADER
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(RIGHT_COLOR)

    # Right bullets
    for i, bullet in enumerate(RIGHT_BULLETS):
        y_pos = 2.5 + (i * 1.2)

        # Bullet marker
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(6.93), Inches(y_pos + 0.12),
            Inches(0.1), Inches(0.1)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = hex_to_rgb(RIGHT_COLOR)
        marker.line.fill.background()

        # Bullet text
        text_box = slide.shapes.add_textbox(
            Inches(7.28), Inches(y_pos),
            Inches(5.5), Inches(1.0)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Save
    output = Path("two-column-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
