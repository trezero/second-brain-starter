#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
# /// layout
# name = "circular-hero-slide"
# purpose = "Central concept, core idea, focus point with surrounding context"
# best_for = [
#     "Core concept with 4-6 related aspects",
#     "Hub-and-spoke relationships",
#     "Central theme with supporting elements",
#     "Ecosystem or component overview",
# ]
# avoid_when = [
#     "Items are sequential (use process/timeline instead)",
#     "More than 6 surrounding items",
#     "Center concept needs long text",
# ]
# min_surrounding_items = 4
# max_surrounding_items = 6
# instructions = [
#     "Center text should be 1-2 words maximum",
#     "Surrounding items are positioned in a circle - keep them short",
#     "4-6 surrounding items work best for visual balance",
#     "All items should relate to the central concept",
#     "Subtitle under center text is optional",
# ]
# ///
"""
LAYOUT: Circular Hero
PURPOSE: Central concept, core idea, focus point with surrounding context

CUSTOMIZE:
- BRAND_* colors and fonts
- CENTER_TEXT: Word or short phrase in the circle
- SURROUNDING_ITEMS: List of related concepts around the circle (4-6 items)
"""

from pathlib import Path
import math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    # === BRAND COLORS (get from brands/{name}/brand-system.md) ===
    BRAND_BG = "REPLACE"              # background color
    BRAND_TEXT = "REPLACE"            # primary text color
    BRAND_TEXT_SECONDARY = "REPLACE"  # secondary/muted text color
    BRAND_ACCENT = "REPLACE"          # accent color (for main circle)
    BRAND_ACCENT_SECONDARY = "REPLACE"  # secondary accent (for surrounding dots)
    BRAND_HEADING_FONT = "REPLACE"    # heading font name
    BRAND_BODY_FONT = "REPLACE"       # body font name

    # === CONTENT (from user request, written in brand voice) ===
    CENTER_TEXT = "REPLACE"           # word or short phrase in the circle
    SUBTITLE = "REPLACE"              # subtitle below center text
    SURROUNDING_ITEMS = [
        "REPLACE",                    # related concepts around the circle (4-6 items)
        "REPLACE",
        "REPLACE",
        "REPLACE",
        "REPLACE",
        "REPLACE",
    ]

    # === SLIDE SETUP ===
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Center coordinates
    center_x = 6.666  # Middle of slide
    center_y = 3.75   # Middle of slide

    # Outer ring (subtle)
    outer_ring = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(center_x - 3.2), Inches(center_y - 3.2),
        Inches(6.4), Inches(6.4)
    )
    outer_ring.fill.background()
    outer_ring.line.color.rgb = hex_to_rgb("313244")
    outer_ring.line.width = Pt(2)

    # Main circle
    circle_size = 3.5
    main_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(center_x - circle_size/2), Inches(center_y - circle_size/2),
        Inches(circle_size), Inches(circle_size)
    )
    main_circle.fill.solid()
    main_circle.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    main_circle.line.fill.background()

    # Center text
    center_box = slide.shapes.add_textbox(
        Inches(center_x - 1.75), Inches(center_y - 0.6),
        Inches(3.5), Inches(1.0)
    )
    tf = center_box.text_frame
    p = tf.paragraphs[0]
    p.text = CENTER_TEXT
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_BG)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle below center
    subtitle_box = slide.shapes.add_textbox(
        Inches(center_x - 2), Inches(center_y + 0.4),
        Inches(4), Inches(0.5)
    )
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = SUBTITLE
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = hex_to_rgb(BRAND_BG)
    p.alignment = PP_ALIGN.CENTER

    # Surrounding items positioned in a circle
    num_items = len(SURROUNDING_ITEMS)
    radius = 3.0  # Distance from center

    for i, item in enumerate(SURROUNDING_ITEMS):
        # Calculate position on circle
        angle = (2 * math.pi * i / num_items) - math.pi / 2  # Start from top
        item_x = center_x + radius * math.cos(angle)
        item_y = center_y + radius * math.sin(angle)

        # Small accent dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(item_x - 0.1), Inches(item_y - 0.1),
            Inches(0.2), Inches(0.2)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY)
        dot.line.fill.background()

        # Item text - positioned based on which side of circle
        text_width = 1.8
        if item_x < center_x - 1:
            text_x = item_x - text_width - 0.2
            align = PP_ALIGN.RIGHT
        elif item_x > center_x + 1:
            text_x = item_x + 0.3
            align = PP_ALIGN.LEFT
        else:
            text_x = item_x - text_width / 2
            align = PP_ALIGN.CENTER

        text_y = item_y - 0.2 if item_y < center_y else item_y - 0.1

        item_box = slide.shapes.add_textbox(
            Inches(text_x), Inches(text_y),
            Inches(text_width), Inches(0.5)
        )
        p = item_box.text_frame.paragraphs[0]
        p.text = item
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
        p.alignment = align

    # Save
    output = Path("circular-hero-slide.pptx")
    prs.save(output)
    print(f"Created {output}")


if __name__ == "__main__":
    main()
