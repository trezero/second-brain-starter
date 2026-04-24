#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///
"""
Generate a single preview deck with ALL cookbook layouts.
Run this to see all available layouts in one presentation.

Usage:
    uv run generate-cookbook-preview.py
"""

from pathlib import Path
import math

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# Brand colors (example brand for preview)
BRAND_BG = "1e1e2e"
BRAND_BG_ALT = "181825"
BRAND_TEXT = "cdd6f4"
BRAND_TEXT_SECONDARY = "bac2de"
BRAND_ACCENT = "fab387"
BRAND_ACCENT_SECONDARY = "89b4fa"
BRAND_ACCENT_TERTIARY = "a6e3a1"
BRAND_CARD_BG = "313244"
BRAND_CODE_BG = "11111b"
BRAND_HEADING_FONT = "JetBrains Mono"
BRAND_BODY_FONT = "Inter"


def add_slide_label(slide, label: str):
    """Add a small label showing which layout this is."""
    label_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(7.1),
        Inches(4), Inches(0.3)
    )
    p = label_box.text_frame.paragraphs[0]
    p.text = f"Layout: {label}"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(10)
    p.font.color.rgb = hex_to_rgb("6c7086")


def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Accent bar at top
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    top_bar.line.fill.background()

    # Headline
    headline_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(11.83), Inches(1.8))
    tf = headline_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Cookbook Preview"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(10.33), Inches(1.0))
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = "All available slide layouts"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(24)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.alignment = PP_ALIGN.CENTER

    add_slide_label(slide, "title-slide")


def create_content_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Left accent bar
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    left_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.0), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Key Principles"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Underline
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(2.5), Inches(0.06))
    underline.fill.solid()
    underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY)
    underline.line.fill.background()

    # Bullets
    bullets = ["Context beats prompting every time", "Iteration is not optional", "Ship fast, fix faster", "The model follows your lead"]
    for i, bullet_text in enumerate(bullets):
        y_pos = 2.0 + (i * 1.1)
        indicator = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y_pos + 0.15), Inches(0.12), Inches(0.12))
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
        indicator.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(1.0), Inches(y_pos), Inches(11.5), Inches(0.9))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet_text
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(22)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    add_slide_label(slide, "content-slide")


def create_two_column_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    LEFT_COLOR = "f38ba8"
    RIGHT_COLOR = "a6e3a1"

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Before vs After"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.63), Inches(1.5), Inches(0.07), Inches(5.5))
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    divider.line.fill.background()

    # Left header
    left_ind = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.7), Inches(0.15), Inches(0.5))
    left_ind.fill.solid()
    left_ind.fill.fore_color.rgb = hex_to_rgb(LEFT_COLOR)
    left_ind.line.fill.background()

    left_header = slide.shapes.add_textbox(Inches(0.85), Inches(1.65), Inches(5.5), Inches(0.6))
    p = left_header.text_frame.paragraphs[0]
    p.text = "The Old Way"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(LEFT_COLOR)

    # Right header
    right_ind = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.93), Inches(1.7), Inches(0.15), Inches(0.5))
    right_ind.fill.solid()
    right_ind.fill.fore_color.rgb = hex_to_rgb(RIGHT_COLOR)
    right_ind.line.fill.background()

    right_header = slide.shapes.add_textbox(Inches(7.28), Inches(1.65), Inches(5.5), Inches(0.6))
    p = right_header.text_frame.paragraphs[0]
    p.text = "The New Way"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(RIGHT_COLOR)

    add_slide_label(slide, "two-column-slide")


def create_quote_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large quote mark
    quote_mark = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2.0), Inches(2.0))
    p = quote_mark.text_frame.paragraphs[0]
    p.text = "\u201C"
    p.font.name = "Georgia"
    p.font.size = Pt(200)
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    # Accent bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.8), Inches(0.12), Inches(2.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    accent_bar.line.fill.background()

    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10.5), Inches(2.5))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The model isn't the bottleneck. You are."
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(36)
    p.font.italic = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10.5), Inches(0.6))
    p = attr_box.text_frame.paragraphs[0]
    p.text = "— Your Name"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    add_slide_label(slide, "quote-slide")


def create_section_break_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large accent block
    accent_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
    accent_block.fill.solid()
    accent_block.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
    accent_block.line.fill.background()

    # Section number
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(3.5), Inches(1.5))
    p = num_box.text_frame.paragraphs[0]
    p.text = "02"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    p.alignment = PP_ALIGN.CENTER

    # Divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(1.5), Inches(0.08), Inches(4.5))
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    divider.line.fill.background()

    # Section title
    title_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.6), Inches(7.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Core Problem"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    add_slide_label(slide, "section-break-slide")


def create_stats_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12.0), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Results Speak"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    stats = [("10x", "Faster Development"), ("85%", "Less Boilerplate"), ("3hrs", "To Production")]
    colors = [BRAND_ACCENT, BRAND_ACCENT_SECONDARY, BRAND_ACCENT_TERTIARY]
    stat_width = 4.0

    for i, ((value, label), color) in enumerate(zip(stats, colors)):
        x_pos = 0.66 + (i * stat_width)

        value_box = slide.shapes.add_textbox(Inches(x_pos), Inches(2.5), Inches(stat_width - 0.3), Inches(1.8))
        p = value_box.text_frame.paragraphs[0]
        p.text = value
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(color)
        p.alignment = PP_ALIGN.CENTER

        label_box = slide.shapes.add_textbox(Inches(x_pos), Inches(4.6), Inches(stat_width - 0.3), Inches(1.0))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.alignment = PP_ALIGN.CENTER

    add_slide_label(slide, "stats-slide")


def create_code_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10.0), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Pattern"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Code block background
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(4.8))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = hex_to_rgb(BRAND_CODE_BG)
    code_bg.line.color.rgb = hex_to_rgb("313244")

    # Terminal dots
    for i, color in enumerate(["f38ba8", "f9e2af", "a6e3a1"]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8 + (i * 0.35)), Inches(1.65), Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb(color)
        dot.line.fill.background()

    # Code content
    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3.8))
    p = code_box.text_frame.paragraphs[0]
    p.text = "def create_agent(config):\n    context = load_context(config.path)\n    return Agent(context=context)"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    add_slide_label(slide, "code-slide")


def create_bold_diagonal_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Diagonal blocks
    block1 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(-1), Inches(3), Inches(8), Inches(5))
    block1.fill.solid()
    block1.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    block1.line.fill.background()

    block2 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(-2), Inches(4.5), Inches(7), Inches(4))
    block2.fill.solid()
    block2.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY)
    block2.line.fill.background()

    # Headline
    headline_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.0), Inches(7.3), Inches(2.0))
    tf = headline_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Break Everything"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.RIGHT

    # Subtext
    subtext_box = slide.shapes.add_textbox(Inches(5.5), Inches(4.2), Inches(7.3), Inches(1.0))
    p = subtext_box.text_frame.paragraphs[0]
    p.text = "Then rebuild it better"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(24)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.alignment = PP_ALIGN.RIGHT

    add_slide_label(slide, "bold-diagonal-slide")


def create_giant_focus_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Background circle
    bg_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(0.5), Inches(6.333), Inches(6.333))
    bg_circle.fill.solid()
    bg_circle.fill.fore_color.rgb = hex_to_rgb("232333")
    bg_circle.line.fill.background()

    # Context above
    above_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(0.6))
    p = above_box.text_frame.paragraphs[0]
    p.text = "NOT INCREMENTAL IMPROVEMENT"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = hex_to_rgb("6c7086")
    p.alignment = PP_ALIGN.CENTER

    # GIANT TEXT
    big_box = slide.shapes.add_textbox(Inches(0), Inches(2.2), Inches(13.333), Inches(3.5))
    p = big_box.text_frame.paragraphs[0]
    p.text = "10x"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(200)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    p.alignment = PP_ALIGN.CENTER

    # Context below
    below_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12.333), Inches(0.6))
    p = below_box.text_frame.paragraphs[0]
    p.text = "Order of magnitude thinking"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    add_slide_label(slide, "giant-focus-slide")


def create_floating_cards_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Three Pillars"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    cards = [("Context", "Feed the model what it needs"), ("Tools", "Give it ability to act"), ("Iteration", "Refine until done")]
    configs = [
        {"x": 0.8, "y": 2.0, "accent": BRAND_ACCENT},
        {"x": 4.5, "y": 2.4, "accent": BRAND_ACCENT_SECONDARY},
        {"x": 8.2, "y": 1.8, "accent": BRAND_ACCENT_TERTIARY},
    ]

    for i, ((title, desc), config) in enumerate(zip(cards, configs)):
        x, y, accent = config["x"], config["y"], config["accent"]

        # Shadow
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(y + 0.15), Inches(4.0), Inches(3.8))
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = hex_to_rgb("11111b")
        shadow.line.fill.background()

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.0), Inches(3.8))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.fill.background()

        # Accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.3), Inches(y + 0.3), Inches(0.8), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(accent)
        bar.line.fill.background()

        # Number
        num = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.6), Inches(1.0), Inches(0.8))
        p = num.text_frame.paragraphs[0]
        p.text = f"0{i+1}"
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(accent)

        # Title
        t = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.5), Inches(3.4), Inches(0.8))
        p = t.text_frame.paragraphs[0]
        p.text = title
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    add_slide_label(slide, "floating-cards-slide")


def create_circular_hero_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    center_x, center_y = 6.666, 3.75

    # Outer ring
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - 3.2), Inches(center_y - 3.2), Inches(6.4), Inches(6.4))
    outer.fill.background()
    outer.line.color.rgb = hex_to_rgb("313244")
    outer.line.width = Pt(2)

    # Main circle
    main = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - 1.75), Inches(center_y - 1.75), Inches(3.5), Inches(3.5))
    main.fill.solid()
    main.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    main.line.fill.background()

    # Center text
    center = slide.shapes.add_textbox(Inches(center_x - 1.5), Inches(center_y - 0.4), Inches(3), Inches(1.0))
    p = center.text_frame.paragraphs[0]
    p.text = "Context"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_BG)
    p.alignment = PP_ALIGN.CENTER

    # Surrounding items
    items = ["Docs", "Examples", "Constraints", "History", "Tools", "Goals"]
    for i, item in enumerate(items):
        angle = (2 * math.pi * i / len(items)) - math.pi / 2
        x = center_x + 2.8 * math.cos(angle)
        y = center_y + 2.8 * math.sin(angle)

        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.08), Inches(y - 0.08), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY)
        dot.line.fill.background()

    add_slide_label(slide, "circular-hero-slide")


def create_chart_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Time Allocation"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Chart
    chart_data = CategoryChartData()
    chart_data.categories = ["Context", "Iteration", "Prompting", "Debugging"]
    chart_data.add_series("Values", [40, 30, 20, 10])

    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.5), Inches(1.8), Inches(7.5), Inches(5.0), chart_data)
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT

    # Insight box
    insight_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(2.5), Inches(4.3), Inches(2.5))
    insight_bg.fill.solid()
    insight_bg.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
    insight_bg.line.fill.background()

    insight_label = slide.shapes.add_textbox(Inches(8.8), Inches(2.7), Inches(3.7), Inches(0.5))
    p = insight_label.text_frame.paragraphs[0]
    p.text = "KEY INSIGHT"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    insight = slide.shapes.add_textbox(Inches(8.8), Inches(3.2), Inches(3.7), Inches(1.5))
    tf = insight.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "40% of success comes from context engineering"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    add_slide_label(slide, "chart-slide")


def create_corner_anchor_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large anchor circle (top-right)
    anchor = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(13.333 - 2.75), Inches(-2.75), Inches(5.5), Inches(5.5))
    anchor.fill.solid()
    anchor.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    anchor.line.fill.background()

    # Secondary circle
    small = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(13.333 - 2.5), Inches(-0.5), Inches(2.5), Inches(2.5))
    small.fill.solid()
    small.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY)
    small.line.fill.background()

    # Headline
    headline = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(7.0), Inches(1.5))
    p = headline.text_frame.paragraphs[0]
    p.text = "Ship Faster"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.6), Inches(2.5), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    line.line.fill.background()

    # Body
    body = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(7.0), Inches(2.0))
    tf = body.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The best code is code running in production, not perfect code in a branch."
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    add_slide_label(slide, "corner-anchor-slide")


def create_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Bottom block
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.0), Inches(13.333), Inches(2.5))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
    bottom.line.fill.background()

    # Top bar
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    top.fill.solid()
    top.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    top.line.fill.background()

    # Headline
    headline = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.33), Inches(1.5))
    p = headline.text_frame.paragraphs[0]
    p.text = "Start Building Today"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # Underline
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.6), Inches(2.33), Inches(0.08))
    underline.fill.solid()
    underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    underline.line.fill.background()

    add_slide_label(slide, "closing-slide")


def create_image_caption_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(1.2))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "See It In Action"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Underline
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.8), Inches(2.0), Inches(0.06))
    underline.fill.solid()
    underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    underline.line.fill.background()

    # Caption
    caption = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(5.5), Inches(3.0))
    tf = caption.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Real-time visibility into your agent's context window and decisions."
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # Image placeholder
    placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(0.8), Inches(6.33), Inches(5.9))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
    placeholder.line.color.rgb = hex_to_rgb("313244")

    # Placeholder text
    img_text = slide.shapes.add_textbox(Inches(8.5), Inches(3.5), Inches(2.0), Inches(0.5))
    p = img_text.text_frame.paragraphs[0]
    p.text = "IMAGE"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(16)
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    p.alignment = PP_ALIGN.CENTER

    add_slide_label(slide, "image-caption-slide")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Generate all layouts
    print("Generating cookbook preview...")
    create_title_slide(prs)
    create_content_slide(prs)
    create_two_column_slide(prs)
    create_quote_slide(prs)
    create_section_break_slide(prs)
    create_stats_slide(prs)
    create_code_slide(prs)
    create_bold_diagonal_slide(prs)
    create_giant_focus_slide(prs)
    create_floating_cards_slide(prs)
    create_circular_hero_slide(prs)
    create_chart_slide(prs)
    create_corner_anchor_slide(prs)
    create_closing_slide(prs)
    create_image_caption_slide(prs)

    output = Path("cookbook-preview.pptx")
    prs.save(output)
    print(f"Created {output} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
